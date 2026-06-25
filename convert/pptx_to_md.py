"""Convert a PowerPoint .pptx file to Markdown (python-pptx)."""
from __future__ import annotations

from pptx import Presentation

from convert import common


def pptx_to_md(in_path):
    prs = Presentation(str(in_path))
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                blocks.append(shape.text.strip())
    return "\n\n".join(blocks) + "\n"


if __name__ == "__main__":
    common.cli(pptx_to_md, in_suffixes={".pptx"}, out_suffix=".md", description=__doc__)
