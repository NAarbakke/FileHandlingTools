import fitz
import pytest


@pytest.fixture
def tiny_pdf(tmp_path):
    """A 1-page PDF with two text lines and one vector 'figure' (a filled rect)."""
    path = tmp_path / "tiny.pdf"
    doc = fitz.open()
    page = doc.new_page(width=300, height=200)
    page.insert_text((20, 40), "Hello world", fontsize=12)
    page.insert_text((20, 80), "Second line here", fontsize=12)
    page.draw_rect(fitz.Rect(20, 110, 120, 170), color=(0, 0, 1), fill=(0.8, 0.8, 1.0))
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def tiny_png(tmp_path):
    """A small PNG image with one line of text, rendered via PyMuPDF (no Pillow)."""
    pdf = fitz.open()
    page = pdf.new_page(width=200, height=120)
    page.insert_text((20, 50), "Handwritten note", fontsize=14)
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
    path = tmp_path / "note.png"
    pix.save(str(path))
    pdf.close()
    return str(path)


@pytest.fixture
def two_page_pdf(tmp_path):
    """A 2-page PDF, one text line per page."""
    path = tmp_path / "two.pdf"
    doc = fitz.open()
    for i in range(2):
        page = doc.new_page(width=300, height=200)
        page.insert_text((20, 40), f"Page {i + 1} text", fontsize=12)
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def tiny_docx(tmp_path):
    """A small .docx with a heading and a body paragraph (python-docx)."""
    from docx import Document
    doc = Document()
    doc.add_heading("Docx Title", level=1)
    doc.add_paragraph("Body paragraph text.")
    path = tmp_path / "tiny.docx"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def tiny_pptx(tmp_path):
    """A small .pptx with a title and a body textbox (python-pptx)."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
    slide.shapes.title.text = "Slide Title"
    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(2))
    box.text_frame.text = "Slide body content."
    path = tmp_path / "tiny.pptx"
    prs.save(str(path))
    return str(path)
